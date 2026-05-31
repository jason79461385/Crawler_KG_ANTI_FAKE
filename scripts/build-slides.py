#!/usr/bin/env python3
"""
Scam Intel Console 作品簡報生成器
- 16:9 寬螢幕
- 深色主題 (與專案 UI 一致)
- 中文 (PingFang TC / Microsoft JhengHei 自動 fallback)
- 4 成員平均分配實作部分
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------- 主題色 ----------
BG = RGBColor(0x0A, 0x14, 0x26)
PANEL = RGBColor(0x0F, 0x1C, 0x33)
PANEL_2 = RGBColor(0x14, 0x23, 0x3F)
BORDER = RGBColor(0x1F, 0x33, 0x54)
TEXT = RGBColor(0xE6, 0xED, 0xF7)
DIM = RGBColor(0x9F, 0xB0, 0xC9)
MUTED = RGBColor(0x6B, 0x7E, 0x9B)
ACCENT = RGBColor(0x38, 0xBD, 0xF8)  # sky-400
ACCENT_2 = RGBColor(0x22, 0xD3, 0xEE)  # cyan-400
ORANGE = RGBColor(0xFB, 0x92, 0x3C)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
YELLOW = RGBColor(0xFA, 0xCC, 0x15)
RED = RGBColor(0xF8, 0x71, 0x71)
PURPLE = RGBColor(0xC0, 0x84, 0xFC)

FONT_TITLE = "PingFang TC"
FONT_BODY = "PingFang TC"
FONT_MONO = "Menlo"

# ---------- 工具函式 ----------
def set_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill=PANEL, line=BORDER, line_w=0.75):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    shp.adjustments[0] = 0.08
    f = shp.fill
    f.solid()
    f.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp

def add_text_box(slide, x, y, w, h, text, *,
                 size=14, color=TEXT, bold=False, font=FONT_BODY,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 line_spacing=1.25):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        r = p.add_run()
        r.text = line
        r.font.name = font
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
    return tb

def add_runs(slide, x, y, w, h, runs, *, align=PP_ALIGN.LEFT,
             anchor=MSO_ANCHOR.TOP, line_spacing=1.25):
    """runs: list of (text, dict(size, color, bold, font))"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    paragraphs = []
    cur = []
    for run in runs:
        if run[0] == "\n":
            paragraphs.append(cur)
            cur = []
        else:
            cur.append(run)
    paragraphs.append(cur)
    for i, para in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        for txt, opts in para:
            r = p.add_run()
            r.text = txt
            r.font.name = opts.get("font", FONT_BODY)
            r.font.size = Pt(opts.get("size", 14))
            r.font.color.rgb = opts.get("color", TEXT)
            r.font.bold = opts.get("bold", False)
    return tb

def add_pill(slide, x, y, text, color, fg=BG):
    width = Inches(max(0.7, 0.13 * len(text) + 0.45))
    height = Inches(0.32)
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, width, height)
    shp.adjustments[0] = 0.5
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = Emu(50000); tf.margin_right = Emu(50000)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(10)
    r.font.bold = True
    r.font.color.rgb = fg
    return shp, width

def add_slide_header(slide, kicker, title):
    add_text_box(slide, Inches(0.55), Inches(0.35),
                 Inches(12), Inches(0.4), kicker.upper(),
                 size=11, color=ACCENT_2, bold=True)
    add_text_box(slide, Inches(0.55), Inches(0.6),
                 Inches(12), Inches(0.8), title,
                 size=30, color=TEXT, bold=True, font=FONT_TITLE)
    # underline accent
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.3),
        Inches(0.6), Emu(38000))
    line.fill.solid(); line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

def add_member_badge(slide, x, y, member_id, color):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, Inches(0.6), Inches(0.6))
    shp.fill.solid()
    shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    tf = shp.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = member_id
    r.font.name = FONT_TITLE
    r.font.size = Pt(20)
    r.font.bold = True
    r.font.color.rgb = BG

def add_footer(slide, idx, total):
    add_text_box(slide, Inches(0.55), Inches(7.0),
                 Inches(6), Inches(0.3),
                 "Scam Intel Console · 防詐情資控制台",
                 size=9, color=MUTED)
    add_text_box(slide, Inches(11.5), Inches(7.0),
                 Inches(1.95), Inches(0.3),
                 f"{idx} / {total}",
                 size=9, color=MUTED, align=PP_ALIGN.RIGHT)

# ---------- 簡報組裝 ----------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

slides_def = []  # 用來事後填頁碼

# ============ Slide 1: 封面 ============
def slide_cover():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    # decorative pill
    add_text_box(s, Inches(0.9), Inches(1.6),
                 Inches(8), Inches(0.5),
                 "FINAL PROJECT REPORT  ·  防詐情資控制台",
                 size=14, color=ACCENT_2, bold=True)
    add_text_box(s, Inches(0.9), Inches(2.2),
                 Inches(12), Inches(1.4),
                 "Scam Intel Console",
                 size=64, color=TEXT, bold=True, font=FONT_TITLE)
    add_text_box(s, Inches(0.9), Inches(3.5),
                 Inches(11.5), Inches(2.4),
                 "整合即時爬蟲 × 知識圖譜 × 規則 + Embedding + LLM 混合分析\n"
                 "× SSRF 防護網址驗證 × Markdown 對話助手 的全端防詐 Prototype",
                 size=22, color=DIM, line_spacing=1.5)
    # team badges
    palette = [ACCENT_2, GREEN, ORANGE, PURPLE]
    for i, c in enumerate(palette):
        add_member_badge(s, Inches(0.9 + i * 0.75), Inches(5.3),
                         "ABCD"[i], c)
    add_text_box(s, Inches(0.9), Inches(6.0),
                 Inches(8), Inches(0.4),
                 "四人協作 · 平均分工",
                 size=14, color=DIM)
    add_text_box(s, Inches(0.9), Inches(6.4),
                 Inches(8), Inches(0.4),
                 "Repo: github.com/jason79461385/Crawler_KG_ANTI_FAKE",
                 size=12, color=MUTED)

slide_cover(); slides_def.append("cover")

# ============ Slide 2: 議程 ============
def slide_agenda():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_slide_header(s, "Agenda", "簡報大綱")
    items = [
        ("01", "專案動機與目標", "為什麼做這個系統"),
        ("02", "系統能力與架構", "六大子系統 + 完整資料流"),
        ("03", "技術棧總覽", "前端 / 後端 / 資料 / 模型"),
        ("04", "四人實作分工", "成員 A / B / C / D 詳述"),
        ("05", "核心演算法", "混合風險評分公式"),
        ("06", "驗證與測試", "型別、API、SSRF、KG"),
        ("07", "Demo 與未來工作", "現場展示 + Roadmap"),
        ("08", "參考資料", "API / 套件 / 學術文獻"),
    ]
    col_w = Inches(6.1)
    for i, (no, t, d) in enumerate(items):
        col = i % 2
        row = i // 2
        x = Inches(0.55 + col * 6.3)
        y = Inches(1.8 + row * 1.25)
        add_rect(s, x, y, col_w, Inches(1.0), fill=PANEL)
        add_text_box(s, x + Inches(0.3), y + Inches(0.15),
                     Inches(0.7), Inches(0.7),
                     no, size=22, color=ACCENT_2, bold=True)
        add_text_box(s, x + Inches(1.1), y + Inches(0.12),
                     Inches(4.8), Inches(0.45),
                     t, size=16, color=TEXT, bold=True)
        add_text_box(s, x + Inches(1.1), y + Inches(0.52),
                     Inches(4.8), Inches(0.4),
                     d, size=11, color=DIM)
slide_agenda(); slides_def.append("agenda")

# ============ Slide 3: 動機 ============
def slide_motivation():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_slide_header(s, "01 · Motivation", "為什麼做這個系統")
    add_text_box(s, Inches(0.55), Inches(1.7),
                 Inches(12.2), Inches(1.0),
                 "詐騙在台灣演化極快:LINE 假投資、假交友、解除分期、165 簡訊、釣魚連結 ……\n"
                 "一般使用者要在訊息瞬間做出「這是不是詐騙」的判斷,缺乏資料支援與即時佐證。",
                 size=15, color=DIM, line_spacing=1.6)
    add_text_box(s, Inches(0.55), Inches(3.0),
                 Inches(12.2), Inches(0.4),
                 "我們希望這個系統能讓使用者:",
                 size=15, color=TEXT, bold=True)
    points = [
        ("即時看見", "案例知識圖譜(誰、用什麼平台、騙什麼)", ACCENT_2),
        ("快速分析", "輸入可疑訊息 → 混合式風險評分 + LLM 解釋", ORANGE),
        ("驗證連結", "可疑網址是否在已知釣魚清單 / 啟發式高風險", GREEN),
        ("自然互動", "浮動聊天助手 + Markdown 排版,跨頁可用", PURPLE),
        ("主動貢獻", "回報遭遇 → 長期累積本地差異化資料", ACCENT),
    ]
    y = Inches(3.55)
    for kicker, body, c in points:
        add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.55),
                 fill=PANEL, line=None)
        dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                  Inches(0.85), y + Inches(0.18),
                                  Inches(0.18), Inches(0.18))
        dot.fill.solid(); dot.fill.fore_color.rgb = c
        dot.line.fill.background()
        add_text_box(s, Inches(1.15), y + Inches(0.1),
                     Inches(2.2), Inches(0.4),
                     kicker, size=13, color=c, bold=True)
        add_text_box(s, Inches(3.4), y + Inches(0.1),
                     Inches(9.2), Inches(0.4),
                     body, size=13, color=TEXT)
        y += Inches(0.65)
slide_motivation(); slides_def.append("motivation")

# ============ Slide 4: 系統能力一覽 ============
def slide_capabilities():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_slide_header(s, "02 · Capabilities", "系統能力一覽")
    stats = [
        ("5", "資料來源", ACCENT_2),
        ("10", "REST API 端點", GREEN),
        ("5", "前端分頁", ORANGE),
        ("3 層", "混合風險評分", PURPLE),
        ("8+", "SSRF 阻擋規則", RED),
        ("7K+", "釣魚 URL 索引", YELLOW),
    ]
    card_w = Inches(2.0)
    gap = Inches(0.05)
    total_w = card_w * 6 + gap * 5
    start_x = (prs.slide_width - total_w) / 2
    y = Inches(1.8)
    for i, (v, k, c) in enumerate(stats):
        x = start_x + (card_w + gap) * i
        add_rect(s, x, y, card_w, Inches(1.4), fill=PANEL_2)
        add_text_box(s, x, y + Inches(0.18),
                     card_w, Inches(0.7), v,
                     size=36, color=c, bold=True, align=PP_ALIGN.CENTER)
        add_text_box(s, x, y + Inches(0.95),
                     card_w, Inches(0.35), k,
                     size=11, color=DIM, align=PP_ALIGN.CENTER)
    # six feature blocks
    blocks = [
        ("即時多來源爬蟲",
         "PTT 5 板 + Google News RSS + 165 全民防騙網 + OpenPhish/PhishTank"),
        ("SQLite 持久化",
         "better-sqlite3 + WAL 模式,URL hash 去重,重啟不流失"),
        ("Neo4j 知識圖譜",
         "MERGE + lastSeenAt 增量同步,Docker Compose 一鍵起"),
        ("混合風險評分",
         "規則 + Embedding 餘弦相似 (×42) + LLM JSON 輸出"),
        ("SSRF 防護",
         "私有 IP / metadata / DNS rebinding + timeout/maxBytes/maxRedirects"),
        ("Markdown 對話助手",
         "浮動圓形 ChatWidget,react-markdown 渲染,localStorage 歷史"),
    ]
    by = Inches(3.55)
    for i, (t, d) in enumerate(blocks):
        col = i % 3
        row = i // 3
        x = Inches(0.55 + col * 4.27)
        y = by + Inches(row * 1.7)
        add_rect(s, x, y, Inches(4.0), Inches(1.55), fill=PANEL)
        add_text_box(s, x + Inches(0.3), y + Inches(0.18),
                     Inches(3.4), Inches(0.4),
                     t, size=14, color=ACCENT_2, bold=True)
        add_text_box(s, x + Inches(0.3), y + Inches(0.6),
                     Inches(3.4), Inches(0.9),
                     d, size=11, color=DIM, line_spacing=1.4)
slide_capabilities(); slides_def.append("capabilities")

# ============ Slide 5: 架構圖 ============
def slide_architecture():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_slide_header(s, "02 · Architecture", "系統架構與資料流")
    art = (
        "┌─────────────────────────────────────────────────────────────────────┐\n"
        "│  React 19 + Vite 7 + Tailwind 4                                       │\n"
        "│  Dashboard · Cases · Verify · Graph · Report  + 浮動 ChatWidget       │\n"
        "└──────────────────────────────┬──────────────────────────────────────┘\n"
        "                               │  /api/* (Vite proxy → :8787)\n"
        "                               ▼\n"
        "┌─────────────────────────────────────────────────────────────────────┐\n"
        "│  Express 5   server/index.ts                                          │\n"
        "│  health · snapshot · feed · graph · crawl · analyze · verify-site      │\n"
        "│  · chat · report · reports · purge-noise                              │\n"
        "└────┬───────────────┬───────────────────────┬─────────────────────────┘\n"
        "     ▼               ▼                       ▼\n"
        " scamEngine.ts    safeFetch.ts          threatIntel.ts\n"
        " (orchestrator)   (SSRF guard)          (PhishTank/OpenPhish/GSB)\n"
        "     │\n"
        "     ├──► SQLite (WAL)               ← posts / user reports / phish urls\n"
        "     └──► Neo4j 5.26 (Docker)        ← MERGE + lastSeenAt 增量同步\n"
        "\n"
        " 外部來源 (任一失敗都 graceful fallback):\n"
        "   PTT(cheerio)  · Google News RSS · 165 CIB_DWS_API · OpenPhish · LLM/Embedding\n"
    )
    tb = add_text_box(s, Inches(0.55), Inches(1.55),
                      Inches(12.2), Inches(5.3),
                      art, size=10, color=ACCENT, font=FONT_MONO,
                      line_spacing=1.15)
    tb.fill.solid()
    # textbox doesn't have fill, draw a panel behind
slide_architecture(); slides_def.append("architecture")

# ============ Slide 6: 技術棧 ============
def slide_techstack():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_slide_header(s, "03 · Tech Stack", "技術棧總覽")
    cats = [
        ("前端", ACCENT_2, [
            "React 19", "Vite 7", "TypeScript 5",
            "Tailwind CSS 4", "vis-network", "react-router 7",
            "react-markdown + remark-gfm", "lucide-react"
        ]),
        ("後端", GREEN, [
            "Node.js 20+", "Express 5", "tsx (TS 直執行)",
            "cors", "dotenv", "cheerio (HTML)",
            "fast-xml-parser (RSS)", "better-sqlite3"
        ]),
        ("資料", ORANGE, [
            "SQLite (WAL)", "Neo4j 5.26 Community", "Docker Compose",
            "in-memory cache + ETag", "URL hash 去重", "lastSeenAt 增量同步"
        ]),
        ("AI / 模型", PURPLE, [
            "Ollama qwen3.5:9b (NCU)", "multilingual-e5-large (embedding)",
            "OpenAI 相容 chat/completions", "response_format: json_object",
            "cosine similarity 評分", "LLM JSON 強制輸出"
        ]),
    ]
    col_w = Inches(3.05)
    gap = Inches(0.15)
    start_x = (prs.slide_width - (col_w * 4 + gap * 3)) / 2
    y = Inches(1.7)
    h = Inches(5.2)
    for i, (title, c, items) in enumerate(cats):
        x = start_x + (col_w + gap) * i
        add_rect(s, x, y, col_w, h, fill=PANEL)
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y,
                                  col_w, Inches(0.05))
        bar.fill.solid(); bar.fill.fore_color.rgb = c
        bar.line.fill.background()
        add_text_box(s, x + Inches(0.25), y + Inches(0.2),
                     col_w - Inches(0.5), Inches(0.5),
                     title, size=18, color=c, bold=True)
        item_y = y + Inches(0.85)
        for it in items:
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                      x + Inches(0.3), item_y + Inches(0.13),
                                      Inches(0.1), Inches(0.1))
            dot.fill.solid(); dot.fill.fore_color.rgb = c
            dot.line.fill.background()
            add_text_box(s, x + Inches(0.5), item_y,
                         col_w - Inches(0.7), Inches(0.4),
                         it, size=11, color=TEXT)
            item_y += Inches(0.5)
slide_techstack(); slides_def.append("techstack")

# ============ Slide 7: 四人分工總覽 ============
def slide_team_overview():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_slide_header(s, "04 · Team Split", "四人實作分工 · 平均分配")
    members = [
        ("A", ACCENT_2, "資料蒐集層", "Data Ingest",
         ["PTT 5 板爬蟲 (cheerio)",
          "Google News RSS (fast-xml-parser)",
          "165 全民防騙網 SPA 反向工程",
          "內容過濾器 BLOCK/ALLOW 雙清單",
          "使用者回報表單與後端寫入"]),
        ("B", GREEN, "持久化 & 安全", "Persistence & Security",
         ["SQLite (better-sqlite3, WAL) 設計",
          "Neo4j MERGE 增量同步 + retry backoff",
          "Docker Compose + db-up 智慧腳本",
          "SSRF 防護 safeFetch",
          "威脅情資 PhishTank/OpenPhish/GSB"]),
        ("C", ORANGE, "智慧分析層", "Intelligence",
         ["scamEngine 評分公式設計",
          "LLM OpenAI 相容 client (Bearer + Basic)",
          "Embedding 語意檢索 + 快取",
          "analyzeMessage / verifySiteUrl",
          "buildKnowledgeGraph 投影"]),
        ("D", PURPLE, "前端 & 體驗", "Frontend & UX",
         ["React Router v7 五頁路由",
          "vis-network KG (全螢幕 + 熱過濾)",
          "ChatWidget 浮動圓形 + Markdown",
          "SnapshotContext 三層快取",
          "API client + 各頁狀態管理"]),
    ]
    card_w = Inches(3.05)
    gap = Inches(0.15)
    start_x = (prs.slide_width - (card_w * 4 + gap * 3)) / 2
    y = Inches(1.65)
    h = Inches(5.4)
    for i, (mid, c, role_tc, role_en, tasks) in enumerate(members):
        x = start_x + (card_w + gap) * i
        add_rect(s, x, y, card_w, h, fill=PANEL)
        add_member_badge(s, x + Inches(0.3), y + Inches(0.3), mid, c)
        add_text_box(s, x + Inches(1.05), y + Inches(0.3),
                     card_w - Inches(1.3), Inches(0.4),
                     role_tc, size=16, color=TEXT, bold=True)
        add_text_box(s, x + Inches(1.05), y + Inches(0.7),
                     card_w - Inches(1.3), Inches(0.35),
                     role_en, size=10, color=DIM)
        ty = y + Inches(1.3)
        for t in tasks:
            dot = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                      x + Inches(0.3), ty + Inches(0.13),
                                      Inches(0.1), Inches(0.1))
            dot.fill.solid(); dot.fill.fore_color.rgb = c
            dot.line.fill.background()
            add_text_box(s, x + Inches(0.5), ty,
                         card_w - Inches(0.7), Inches(0.7),
                         t, size=11, color=TEXT, line_spacing=1.3)
            ty += Inches(0.75)
slide_team_overview(); slides_def.append("team_overview")

# ============ 通用「成員頁」模板 ============
def member_intro_slide(mid, color, title_tc, title_en, summary_lines):
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    # 大徽章
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(0.6), Inches(0.5),
                                Inches(1.0), Inches(1.0))
    badge.fill.solid(); badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = mid; r.font.name = FONT_TITLE
    r.font.size = Pt(40); r.font.bold = True; r.font.color.rgb = BG
    add_text_box(s, Inches(1.85), Inches(0.5),
                 Inches(10), Inches(0.45),
                 f"成員 {mid} · {title_en}".upper(),
                 size=12, color=color, bold=True)
    add_text_box(s, Inches(1.85), Inches(0.85),
                 Inches(11), Inches(0.8),
                 title_tc, size=30, color=TEXT, bold=True, font=FONT_TITLE)
    add_text_box(s, Inches(0.55), Inches(1.95),
                 Inches(12.2), Inches(0.6),
                 "本人負責的範圍 · 與其他成員的銜接點",
                 size=12, color=DIM, bold=True)
    # 三大欄(summary blocks)
    cards = summary_lines  # list of (kicker, body)
    cw = Inches(4.0)
    gap = Inches(0.15)
    start_x = (prs.slide_width - (cw * 3 + gap * 2)) / 2
    y = Inches(2.5)
    h = Inches(2.6)
    for i, (k, b) in enumerate(cards):
        x = start_x + (cw + gap) * i
        add_rect(s, x, y, cw, h, fill=PANEL)
        add_text_box(s, x + Inches(0.3), y + Inches(0.25),
                     cw - Inches(0.6), Inches(0.5),
                     k, size=14, color=color, bold=True)
        add_text_box(s, x + Inches(0.3), y + Inches(0.85),
                     cw - Inches(0.6), h - Inches(1.0),
                     b, size=11, color=TEXT, line_spacing=1.5)
    # 底部 chips
    add_text_box(s, Inches(0.55), Inches(5.4),
                 Inches(12.2), Inches(0.4),
                 "本成員主要交付的檔案", size=12, color=DIM, bold=True)
    return s

def detail_slide(mid, color, kicker, title, sections):
    """sections: list of (header, body_lines, optional_code)"""
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    # mini badge
    badge = s.shapes.add_shape(MSO_SHAPE.OVAL,
                                Inches(0.55), Inches(0.4),
                                Inches(0.55), Inches(0.55))
    badge.fill.solid(); badge.fill.fore_color.rgb = color
    badge.line.fill.background()
    tf = badge.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = mid; r.font.name = FONT_TITLE
    r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = BG
    add_text_box(s, Inches(1.3), Inches(0.35),
                 Inches(11), Inches(0.4),
                 kicker.upper(),
                 size=11, color=color, bold=True)
    add_text_box(s, Inches(1.3), Inches(0.6),
                 Inches(11.5), Inches(0.8),
                 title, size=25, color=TEXT, bold=True, font=FONT_TITLE)
    line = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1.3), Inches(1.3),
        Inches(0.5), Emu(38000))
    line.fill.solid(); line.fill.fore_color.rgb = color
    line.line.fill.background()
    return s

# ============ 成員 A 系列 ============
def slides_member_a():
    color = ACCENT_2
    member_intro_slide(
        "A", color, "資料蒐集層", "Data Ingest",
        [("負責範圍",
          "把各種異質來源的詐騙線索抓回來,正規化成統一格式 DemoPost,"
          "並過濾掉政治新聞 / PSA 等噪聲。"),
         ("關鍵挑戰",
          "165 全民防騙網是 Angular SPA,純爬 HTML 抓不到;\n"
          "新聞標題容易誤觸 (例:政治人物姓名含「詐騙」)。"),
         ("交付產物",
          "server/sources/ptt.ts\nserver/sources/googleNews.ts\n"
          "server/sources/dashboard165.ts\nserver/lib/contentFilter.ts")])

    # A-1: PTT
    s = detail_slide("A", color, "Member A · Detail 1",
                     "PTT 5 板爬蟲(cheerio)",
                     [])
    add_runs(s, Inches(1.3), Inches(1.6),
             Inches(7.5), Inches(0.45),
             [("板別:", {"size": 12, "bold": True, "color": DIM}),
              (" Gossiping · e-shopping · part-time · Salary · MobilePay",
               {"size": 12})])
    add_text_box(s, Inches(1.3), Inches(2.05),
                 Inches(7.5), Inches(0.4),
                 "處理流程", size=13, color=color, bold=True)
    steps = [
        "① 帶 over18 cookie 抓 board index,cheerio 解析 .r-ent",
        "② 用 regex 過濾標題:詐騙|被騙|騙錢|假網拍|釣魚",
        "③ 進文章頁抓 #main-content innerText,移除推文/簽名",
        "④ 進 postUtils.inferEntities 抽 entity (line/匯款/平台...)",
        "⑤ inferScamType 推斷 scamType,組成 DemoPost 統一物件",
    ]
    yy = Inches(2.5)
    for st in steps:
        add_text_box(s, Inches(1.3), yy, Inches(7.5), Inches(0.45),
                     st, size=11, color=TEXT, line_spacing=1.5)
        yy += Inches(0.5)
    # code snippet on right
    add_rect(s, Inches(9.1), Inches(1.55),
             Inches(3.65), Inches(5.4), fill=PANEL_2)
    add_text_box(s, Inches(9.3), Inches(1.7),
                 Inches(3.5), Inches(0.4),
                 "關鍵程式碼", size=11, color=color, bold=True)
    snippet = (
        "const $ = cheerio.load(html);\n"
        "$(\".r-ent .title a\").each((_, a) => {\n"
        "  const title = $(a).text();\n"
        "  if (!SCAM_RE.test(title)) return;\n"
        "  const href = $(a).attr(\"href\");\n"
        "  candidates.push({ title, href });\n"
        "});\n"
        "\n"
        "// 個別抓內文\n"
        "const post = await fetch(url, {\n"
        "  headers: { Cookie: \"over18=1\" }\n"
        "});\n"
    )
    add_text_box(s, Inches(9.3), Inches(2.1),
                 Inches(3.35), Inches(4.8),
                 snippet, size=9, color=ACCENT, font=FONT_MONO,
                 line_spacing=1.3)

    # A-2: Google News + 165
    s = detail_slide("A", color, "Member A · Detail 2",
                     "Google News RSS + 165 SPA 反向工程",
                     [])
    add_rect(s, Inches(1.3), Inches(1.6),
             Inches(5.5), Inches(5.4), fill=PANEL)
    add_text_box(s, Inches(1.5), Inches(1.75),
                 Inches(5.0), Inches(0.4),
                 "Google News RSS", size=14, color=color, bold=True)
    body1 = (
        "• 12 組詐騙關鍵字(投資詐騙、解除分期、假交友 ...)\n"
        "• 每組 query 走 news.google.com/rss/search?q=\n"
        "• fast-xml-parser 解析 <item>\n"
        "• 標題 + description 餵內容過濾器再決定要不要收\n"
        "• publishedAt = pubDate;url = link"
    )
    add_text_box(s, Inches(1.5), Inches(2.2),
                 Inches(5.1), Inches(4.6),
                 body1, size=11, color=TEXT, line_spacing=1.7)

    add_rect(s, Inches(7.0), Inches(1.6),
             Inches(5.75), Inches(5.4), fill=PANEL)
    add_text_box(s, Inches(7.2), Inches(1.75),
                 Inches(5.0), Inches(0.4),
                 "165 全民防騙網(SPA 反向工程)", size=14, color=color, bold=True)
    body2 = (
        "問題:165dashboard.tw 是 Angular SPA,純 HTML 沒列表\n\n"
        "解法:逆向 JS bundle 找到真實 API\n"
        "  POST  CIB_DWS_API/api/CaseSummary/GetCaseSummaryList\n"
        "  Body: { PageIndex, NumberOfPerPage, SortOrderInfos }\n"
        "  Resp: 160,999 筆案例\n\n"
        "踩雷:所有案例 url 都指向 /categorized-fraud,\n"
        "       URL hash 去重會把全部折成 1 筆\n"
        "修正:改用 /case-detail/{Id} 當作每筆獨立 URL"
    )
    add_text_box(s, Inches(7.2), Inches(2.2),
                 Inches(5.4), Inches(4.6),
                 body2, size=11, color=TEXT, line_spacing=1.7)

    # A-3: 內容過濾
    s = detail_slide("A", color, "Member A · Detail 3",
                     "內容過濾器:BLOCK + ALLOW 雙清單",
                     [])
    add_text_box(s, Inches(1.3), Inches(1.6),
                 Inches(11), Inches(0.5),
                 "問題:Google News 抓回來的「詐騙」常常是政治新聞 / PSA 宣導 / 統計排行,"
                 "汙染 KG 與 RAG evidence。",
                 size=12, color=DIM, line_spacing=1.5)
    # BLOCK
    add_rect(s, Inches(1.3), Inches(2.5),
             Inches(5.5), Inches(4.0), fill=PANEL)
    add_text_box(s, Inches(1.5), Inches(2.6),
                 Inches(5.1), Inches(0.4),
                 "BLOCK_PATTERNS", size=13, color=RED, bold=True)
    block_items = [
        "/政黨|政治人物|選舉|候選人/",
        "/防範宣導|宣導活動|防詐宣導/",
        "/排行|統計|報告/",
        "/演藝|綜藝|戲劇/",
    ]
    yy = Inches(3.05)
    for it in block_items:
        add_text_box(s, Inches(1.6), yy, Inches(5), Inches(0.4),
                     "• " + it, size=11, color=TEXT, font=FONT_MONO)
        yy += Inches(0.5)
    add_text_box(s, Inches(1.5), Inches(5.4),
                 Inches(5.1), Inches(1.0),
                 "命中即丟棄 (BLOCK)", size=11, color=DIM)

    # ALLOW
    add_rect(s, Inches(7.0), Inches(2.5),
             Inches(5.75), Inches(4.0), fill=PANEL)
    add_text_box(s, Inches(7.2), Inches(2.6),
                 Inches(5.3), Inches(0.4),
                 "ALLOW_PATTERNS", size=13, color=GREEN, bold=True)
    allow_items = [
        "/被騙|遭騙|誤信/",
        "/損失金額|匯款|轉帳/",
        "/接觸管道|LINE|Telegram/",
        "/集團車手|車手|提款卡/",
        "/165 專線|報案/",
    ]
    yy = Inches(3.05)
    for it in allow_items:
        add_text_box(s, Inches(7.3), yy, Inches(5.3), Inches(0.4),
                     "• " + it, size=11, color=TEXT, font=FONT_MONO)
        yy += Inches(0.5)
    add_text_box(s, Inches(7.2), Inches(5.7),
                 Inches(5.3), Inches(1.0),
                 "命中強信號才放行 (ALLOW)", size=11, color=DIM)
slides_member_a()
slides_def.extend(["a1", "a2", "a3", "a4"])

# ============ 成員 B 系列 ============
def slides_member_b():
    color = GREEN
    member_intro_slide(
        "B", color, "持久化 & 安全", "Persistence & Security",
        [("負責範圍",
          "SQLite + Neo4j 雙層持久化、Docker Compose 一鍵啟動、"
          "SSRF 防護以及威脅情資 feed 整合。"),
         ("關鍵挑戰",
          "Neo4j 原本每次 MATCH DETACH DELETE 全砍 → 改 MERGE 增量;\n"
          "verify-site 直接 fetch 使用者 URL → SSRF 高風險。"),
         ("交付產物",
          "server/lib/db.ts\nserver/lib/neo4j.ts\n"
          "server/lib/safeFetch.ts\nserver/lib/threatIntel.ts\n"
          "docker-compose.yml\nscripts/db-up.ts")])

    # B-1: SQLite
    s = detail_slide("B", color, "Member B · Detail 1",
                     "SQLite 持久化(better-sqlite3, WAL)",
                     [])
    add_text_box(s, Inches(1.3), Inches(1.6),
                 Inches(7.5), Inches(0.4),
                 "Schema 設計", size=14, color=color, bold=True)
    schema = (
        "posts            ← id, source, board, title, content, url,\n"
        "                   url_hash, title_hash, scam_type, entities,\n"
        "                   published_at, created_at\n"
        "\n"
        "user_reports     ← id, message, reporter_hint, suspected_url,\n"
        "                   risk_level, risk_score, matched_keywords,\n"
        "                   created_at, status\n"
        "\n"
        "phish_urls       ← url, source, fetched_at\n"
        "\n"
        "source_status    ← source, last_run_at, ok, error_msg, count"
    )
    add_text_box(s, Inches(1.3), Inches(2.1),
                 Inches(7.5), Inches(4.5),
                 schema, size=11, color=TEXT, font=FONT_MONO,
                 line_spacing=1.4)
    add_rect(s, Inches(9.1), Inches(1.55),
             Inches(3.65), Inches(5.4), fill=PANEL_2)
    add_text_box(s, Inches(9.3), Inches(1.7),
                 Inches(3.5), Inches(0.4),
                 "關鍵設計", size=11, color=color, bold=True)
    notes = (
        "✓ WAL 模式 (PRAGMA journal_mode=WAL)\n"
        "✓ URL hash + title hash 雙重去重\n"
        "✓ upsertPost 採 INSERT OR REPLACE\n"
        "✓ 檔案位置 data/scam-intel.db\n"
        "  (.gitignore 排除)\n"
        "✓ 重啟後完整保留 16+ 案例與\n"
        "   使用者回報"
    )
    add_text_box(s, Inches(9.3), Inches(2.1),
                 Inches(3.35), Inches(5.0),
                 notes, size=10, color=TEXT, line_spacing=1.6)

    # B-2: Neo4j MERGE
    s = detail_slide("B", color, "Member B · Detail 2",
                     "Neo4j MERGE 增量同步",
                     [])
    add_rect(s, Inches(1.3), Inches(1.6),
             Inches(5.5), Inches(2.5), fill=PANEL)
    add_text_box(s, Inches(1.5), Inches(1.75),
                 Inches(5.1), Inches(0.4),
                 "舊版策略", size=13, color=RED, bold=True)
    old = (
        "MATCH (n) DETACH DELETE n  -- 全砍重建\n"
        "\n問題:\n"
        "• 資料一多就很慢\n"
        "• 破壞手動加在 Neo4j 上的關係\n"
        "• 連線一失敗就永久禁用,要重啟"
    )
    add_text_box(s, Inches(1.5), Inches(2.2),
                 Inches(5.1), Inches(1.8),
                 old, size=11, color=TEXT, line_spacing=1.5,
                 font=FONT_MONO)

    add_rect(s, Inches(7.0), Inches(1.6),
             Inches(5.75), Inches(2.5), fill=PANEL)
    add_text_box(s, Inches(7.2), Inches(1.75),
                 Inches(5.3), Inches(0.4),
                 "新版策略", size=13, color=GREEN, bold=True)
    new = (
        "MERGE (p:Post {id: $id})\n"
        "SET p.lastSeenAt = $now\n"
        "MERGE (p)-[r:MENTIONS]->(e:Entity)\n"
        "\n清理:\n"
        "WHERE p.lastSeenAt < $cutoff (7d)\n"
        "DETACH DELETE p"
    )
    add_text_box(s, Inches(7.2), Inches(2.2),
                 Inches(5.3), Inches(1.8),
                 new, size=11, color=TEXT, line_spacing=1.5,
                 font=FONT_MONO)

    add_rect(s, Inches(1.3), Inches(4.3),
             Inches(11.45), Inches(2.7), fill=PANEL_2)
    add_text_box(s, Inches(1.5), Inches(4.45),
                 Inches(11), Inches(0.4),
                 "其他可靠性設計", size=13, color=color, bold=True)
    reliability = (
        "✓ connectionTimeout=2000ms,避免本機沒 Neo4j 第一次卡 30s\n"
        "✓ probeNeo4jOnStartup() — 啟動階段就吃掉連線成本,首次 /api/graph < 200ms\n"
        "✓ 失敗 60 秒指數退避重試,而非永久禁用\n"
        "✓ STALE_HOURS = 7 天清理舊節點\n"
        "✓ Docker Compose + healthcheck + db-up.ts 智慧啟動(熱啟動 1.3s)"
    )
    add_text_box(s, Inches(1.5), Inches(4.9),
                 Inches(11), Inches(2.0),
                 reliability, size=11, color=TEXT, line_spacing=1.8)

    # B-3: SSRF
    s = detail_slide("B", color, "Member B · Detail 3",
                     "SSRF 多層防護(safeFetch)",
                     [])
    add_text_box(s, Inches(1.3), Inches(1.6),
                 Inches(11.5), Inches(0.5),
                 "問題:verify-site 直接 fetch 使用者輸入 URL → "
                 "可被當跳板探測內網 / metadata 服務",
                 size=12, color=DIM, line_spacing=1.5)
    # 三欄阻擋類別
    cols = [
        ("阻擋 IP",
         "127.0.0.1 ::1\n"
         "10.0.0.0/8\n172.16.0.0/12\n192.168.0.0/16\n"
         "169.254.0.0/16 (link-local)\n"
         "100.64.0.0/10 (CGNAT)\n"
         "fc00::/7 (ULA IPv6)\n"
         "fe80::/10 (IPv6 link-local)"),
        ("阻擋 Hostname",
         "localhost\nmetadata.google.internal\n"
         "*.local / *.internal\n\n"
         "+ DNS rebinding 防護:\n"
         "  resolve 後再次驗證\n"
         "  IP 不在白名單範圍"),
        ("流量限制",
         "timeout: 5s (預設)\n"
         "maxBytes: 256KB (truncate)\n"
         "maxRedirects: 3 (拋錯)\n"
         "Accept: text/html only\n\n"
         "✓ 三道閘任一觸發即\n"
         "  ssrfBlocked=true 回報"),
    ]
    cw = Inches(3.7)
    gap = Inches(0.15)
    sx = Inches(1.3)
    y = Inches(2.5)
    for i, (k, body) in enumerate(cols):
        x = sx + (cw + gap) * i
        add_rect(s, x, y, cw, Inches(4.4), fill=PANEL)
        add_text_box(s, x + Inches(0.25), y + Inches(0.2),
                     cw - Inches(0.5), Inches(0.4),
                     k, size=13, color=color, bold=True)
        add_text_box(s, x + Inches(0.25), y + Inches(0.7),
                     cw - Inches(0.5), Inches(3.6),
                     body, size=11, color=TEXT,
                     font=FONT_MONO, line_spacing=1.5)
slides_member_b()
slides_def.extend(["b1", "b2", "b3", "b4"])

# ============ 成員 C 系列 ============
def slides_member_c():
    color = ORANGE
    member_intro_slide(
        "C", color, "智慧分析層", "Intelligence",
        [("負責範圍",
          "把抓回來的案例變成「對話風險評分 + 警示」。"
          "整合規則、Embedding、LLM 三層,任一層失效都能 fallback。"),
         ("關鍵挑戰",
          "規則太死、純 LLM 太昂貴又會幻覺。\n"
          "用混合式評分:Embedding 餵相似案例,LLM 只負責文字解釋。"),
         ("交付產物",
          "server/lib/scamEngine.ts\nserver/lib/llm.ts\n"
          "server/lib/postUtils.ts\nbuildKnowledgeGraph()")])

    # C-1: 評分公式
    s = detail_slide("C", color, "Member C · Detail 1",
                     "混合風險評分公式",
                     [])
    add_rect(s, Inches(1.3), Inches(1.6),
             Inches(11.45), Inches(2.4), fill=PANEL_2)
    add_text_box(s, Inches(1.5), Inches(1.75),
                 Inches(11), Inches(0.4),
                 "總分", size=13, color=color, bold=True)
    formula1 = (
        "score = min(96,\n"
        "  10\n"
        "  + matchedKeywords.length * 8\n"
        "  + matchedEntities.length * 10\n"
        "  + scoredEvidence.length * 14\n"
        ")"
    )
    add_text_box(s, Inches(1.5), Inches(2.2),
                 Inches(11), Inches(1.8),
                 formula1, size=12, color=ACCENT, font=FONT_MONO,
                 line_spacing=1.5)

    add_rect(s, Inches(1.3), Inches(4.15),
             Inches(11.45), Inches(2.85), fill=PANEL_2)
    add_text_box(s, Inches(1.5), Inches(4.3),
                 Inches(11), Inches(0.4),
                 "單篇 evidence 分數(三層加權)", size=13, color=color, bold=True)
    formula2 = (
        "evidenceScore =\n"
        "  overlapEntities  * 24      ← 規則命中的實體\n"
        "  + keywordHits    * 14      ← 關鍵字\n"
        "  + cosineSimilarity * 42    ← Embedding (僅當設定時)\n"
        "\n"
        "level:  score >= 70 → high\n"
        "        score >= 40 → medium\n"
        "        otherwise   → low"
    )
    add_text_box(s, Inches(1.5), Inches(4.75),
                 Inches(11), Inches(2.2),
                 formula2, size=11, color=ACCENT, font=FONT_MONO,
                 line_spacing=1.5)

    # C-2: LLM 整合
    s = detail_slide("C", color, "Member C · Detail 2",
                     "LLM 整合(OpenAI 相容 client)",
                     [])
    cards = [
        ("彈性 endpoint",
         "支援任何 OpenAI 相容 /v1/chat/completions:\n"
         "• Ollama (本地)\n• vLLM (NCU GPU)\n"
         "• OpenAI / Azure\n• 自架 LLM gateway"),
        ("雙模式驗證",
         "buildApiRequest 自動判別:\n"
         "• Bearer token (Authorization header)\n"
         "• URL 內嵌 Basic auth\n"
         "  https://user:pass@host/v1"),
        ("強制 JSON 輸出",
         "response_format: json_object\n"
         "schema: { summary, actions[] }\n"
         "失敗自動 fallback 到 buildAlert\n"
         "(deterministic 規則模板)"),
        ("Chat 走 Markdown",
         "/api/chat 的 system prompt 強制\n"
         "輸出 Markdown:標題、條列、表格、\n"
         "粗體、code、警告符號。\n"
         "前端 react-markdown 渲染。"),
    ]
    cw = Inches(5.6)
    gap = Inches(0.2)
    sx = Inches(1.3)
    for i, (k, b) in enumerate(cards):
        col = i % 2
        row = i // 2
        x = sx + col * (cw + gap)
        y = Inches(1.65 + row * 2.7)
        add_rect(s, x, y, cw, Inches(2.5), fill=PANEL)
        add_text_box(s, x + Inches(0.3), y + Inches(0.2),
                     cw - Inches(0.6), Inches(0.4),
                     k, size=14, color=color, bold=True)
        add_text_box(s, x + Inches(0.3), y + Inches(0.7),
                     cw - Inches(0.6), Inches(1.7),
                     b, size=11, color=TEXT, line_spacing=1.6)

    # C-3: Embedding + KG
    s = detail_slide("C", color, "Member C · Detail 3",
                     "Embedding 檢索 + 知識圖譜投影",
                     [])
    add_rect(s, Inches(1.3), Inches(1.6),
             Inches(5.5), Inches(5.4), fill=PANEL)
    add_text_box(s, Inches(1.5), Inches(1.75),
                 Inches(5), Inches(0.4),
                 "Embedding 流程", size=14, color=color, bold=True)
    emb = (
        "1. crawlLiveSources 結束後,\n"
        "   為每篇新 post 預先計算 embedding\n"
        "2. 結果寫入 embeddingCache(per-post id)\n"
        "3. analyzeMessage 時:\n"
        "   • 對 user message 取 embedding\n"
        "   • 與所有 cached embedding 算 cosine\n"
        "   • 依分數排序取 top-N 當 evidence\n"
        "4. 採 OpenAI 相容 /v1/embeddings:\n"
        "   multilingual-e5-large (1024 維)"
    )
    add_text_box(s, Inches(1.5), Inches(2.2),
                 Inches(5.1), Inches(4.6),
                 emb, size=11, color=TEXT, line_spacing=1.7)

    add_rect(s, Inches(7.0), Inches(1.6),
             Inches(5.75), Inches(5.4), fill=PANEL)
    add_text_box(s, Inches(7.2), Inches(1.75),
                 Inches(5), Inches(0.4),
                 "KG 投影邏輯", size=14, color=color, bold=True)
    kg = (
        "buildKnowledgeGraph(posts, maxNodes)\n\n"
        "節點:\n"
        "  post:{id}  ← 案例,group='post'\n"
        "  {type}:{value}  ← 實體\n"
        "    line/匯款/平台/帳戶/keyword\n\n"
        "邊:\n"
        "  (post) -[:MENTIONS]-> (entity)\n\n"
        "權重:\n"
        "  post.weight = entities.length\n"
        "  entity.weight = 出現次數總和\n"
        "  最後依 weight 取 top-N"
    )
    add_text_box(s, Inches(7.2), Inches(2.2),
                 Inches(5.4), Inches(4.6),
                 kg, size=11, color=TEXT, line_spacing=1.6,
                 font=FONT_MONO)
slides_member_c()
slides_def.extend(["c1", "c2", "c3", "c4"])

# ============ 成員 D 系列 ============
def slides_member_d():
    color = PURPLE
    member_intro_slide(
        "D", color, "前端 & 體驗", "Frontend & UX",
        [("負責範圍",
          "React 19 + Vite 7 全端體驗。從 5 個獨立路由頁、KG 視覺化、"
          "浮動聊天到三層 cache。"),
         ("關鍵挑戰",
          "原版單頁 App.tsx 已 692 行;KG 第一次載入 6.7s;\n"
          "節點點擊後 label 看不清楚;chip 過濾要可熱載。"),
         ("交付產物",
          "src/pages/* (5 頁)\nsrc/components/KnowledgeGraphPanel.tsx\n"
          "src/components/ChatWidget.tsx\nsrc/context/SnapshotContext.tsx\n"
          "src/api/client.ts\nsrc/components/Markdown.tsx")])

    # D-1: 路由與 5 頁
    s = detail_slide("D", color, "Member D · Detail 1",
                     "React Router v7:5 頁拆分",
                     [])
    pages = [
        ("/", "Dashboard 主儀表板", "多來源狀態 + scamType 摘要 + 對話分析"),
        ("/cases", "案例列表", "分頁 12 筆/頁 + 來源 chip 過濾"),
        ("/verify", "網站驗證", "6 個 signal chip,SSRF 阻擋會顯示原因"),
        ("/graph", "知識圖譜", "vis-network + 全螢幕 + 熱過濾 + 跳轉來源"),
        ("/report", "使用者回報", "送出即時分析 + 列出近期 20 筆"),
    ]
    y = Inches(1.6)
    for path, name, desc in pages:
        add_rect(s, Inches(1.3), y, Inches(11.45), Inches(0.9), fill=PANEL)
        add_text_box(s, Inches(1.55), y + Inches(0.2),
                     Inches(1.8), Inches(0.5),
                     path, size=13, color=color, bold=True, font=FONT_MONO)
        add_text_box(s, Inches(3.5), y + Inches(0.15),
                     Inches(3.2), Inches(0.4),
                     name, size=14, color=TEXT, bold=True)
        add_text_box(s, Inches(3.5), y + Inches(0.5),
                     Inches(8.5), Inches(0.4),
                     desc, size=11, color=DIM)
        y += Inches(1.0)

    # D-2: KG visualisation
    s = detail_slide("D", color, "Member D · Detail 2",
                     "KG 視覺化:vis-network + 熱過濾 + 全螢幕",
                     [])
    items = [
        ("vis-network DataSet 熱過濾",
         "類別 chip 切換時:DataSet.update + hidden flag + physics 重啟 700ms,"
         "無需刷新頁面;新增分類後節點立即顯示。"),
        ("全螢幕真填滿視窗",
         "fullscreen=true 時改 fixed inset-0 flex 縱向:工具列在上、KG 100% 填滿;"
         "側邊詳情變右側滑出抽屜,點節點自動展開,Esc 離開。"),
        ("案例節點一鍵跳轉來源",
         "GraphNode 新增 url/source/scamType/publishedAt 欄位;NodeDetails 在 post 節點"
         "顯示來源/類型/時間 + 「前往原始來源」藍色按鈕(target=_blank)。"),
        ("節點配色與形狀",
         "shape:dot,label 在外;post 直徑 18-34,entity 8-20。"
         "Stroke 3px 描邊讓深底也能讀標題。Type breakdown 用 group 計算才對。"),
    ]
    y = Inches(1.6)
    for k, b in items:
        add_rect(s, Inches(1.3), y, Inches(11.45), Inches(1.25), fill=PANEL)
        add_text_box(s, Inches(1.55), y + Inches(0.15),
                     Inches(11), Inches(0.4),
                     k, size=13, color=color, bold=True)
        add_text_box(s, Inches(1.55), y + Inches(0.55),
                     Inches(11), Inches(0.7),
                     b, size=11, color=TEXT, line_spacing=1.5)
        y += Inches(1.35)

    # D-3: ChatWidget + SnapshotContext
    s = detail_slide("D", color, "Member D · Detail 3",
                     "ChatWidget + 三層快取",
                     [])
    add_rect(s, Inches(1.3), Inches(1.6),
             Inches(5.5), Inches(5.4), fill=PANEL)
    add_text_box(s, Inches(1.5), Inches(1.75),
                 Inches(5), Inches(0.4),
                 "浮動 ChatWidget", size=14, color=color, bold=True)
    chat = (
        "• 右下角圓形按鈕(漸層 cyan→sky)\n"
        "• 點擊開合,展開 380×560 對話窗\n"
        "• Enter 送出 / Shift+Enter 換行\n"
        "• 對話歷史寫入 localStorage\n"
        "  key: scam-intel-chat-history\n"
        "• 助理回覆走 react-markdown + remark-gfm\n"
        "• 支援 粗體 / code / 表格 / 條列 / 連結\n"
        "• LLM vs 規則模式 header 提示"
    )
    add_text_box(s, Inches(1.5), Inches(2.2),
                 Inches(5.1), Inches(4.6),
                 chat, size=11, color=TEXT, line_spacing=1.7)

    add_rect(s, Inches(7.0), Inches(1.6),
             Inches(5.75), Inches(5.4), fill=PANEL)
    add_text_box(s, Inches(7.2), Inches(1.75),
                 Inches(5), Inches(0.4),
                 "三層 KG 快取", size=14, color=color, bold=True)
    cache = (
        "① 後端 in-memory cache (per limit)\n"
        "    TTL 30s,etag = nodes+edges+posts\n"
        "\n"
        "② HTTP ETag + Cache-Control\n"
        "    304 命中時不重算 graph\n"
        "\n"
        "③ 前端 SnapshotProvider prefetch\n"
        "    /api/snapshot 後立刻背景抓 graph\n"
        "    GraphPage useState 初值即用 cache\n"
        "\n"
        "效果:首次 /graph 6.7s → < 200ms"
    )
    add_text_box(s, Inches(7.2), Inches(2.2),
                 Inches(5.4), Inches(4.6),
                 cache, size=11, color=TEXT, line_spacing=1.7)
slides_member_d()
slides_def.extend(["d1", "d2", "d3", "d4"])

# ============ API 規格 ============
def slide_api():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_slide_header(s, "05 · API Spec", "後端 REST 介面 (10 endpoints)")
    rows = [
        ("GET",  "/api/health",      "健檢 + 子系統狀態",         ACCENT),
        ("GET",  "/api/snapshot",    "多來源狀態 + DB stats",     ACCENT),
        ("GET",  "/api/feed",        "分頁案例清單",              ACCENT),
        ("GET",  "/api/graph",       "KG 節點與邊 (ETag)",        ACCENT),
        ("POST", "/api/crawl",       "觸發即時同步",              GREEN),
        ("POST", "/api/analyze",     "對話風險分析",              GREEN),
        ("POST", "/api/verify-site", "網址啟發式 + SSRF 防護",    GREEN),
        ("POST", "/api/chat",        "ChatWidget 後端 (Markdown)", GREEN),
        ("POST", "/api/report",      "使用者回報 + 自動分析",     GREEN),
        ("GET",  "/api/reports",     "近期使用者回報清單",        ACCENT),
    ]
    add_rect(s, Inches(0.55), Inches(1.55),
             Inches(12.2), Inches(0.5), fill=PANEL_2)
    add_text_box(s, Inches(0.85), Inches(1.62),
                 Inches(1), Inches(0.4),
                 "METHOD", size=11, color=color_for_method("HEAD"), bold=True)
    add_text_box(s, Inches(2.0), Inches(1.62),
                 Inches(4), Inches(0.4),
                 "PATH", size=11, color=DIM, bold=True)
    add_text_box(s, Inches(6.5), Inches(1.62),
                 Inches(6), Inches(0.4),
                 "用途", size=11, color=DIM, bold=True)
    y = Inches(2.1)
    for method, path, desc, c in rows:
        bg = PANEL if (rows.index((method, path, desc, c)) % 2 == 0) else PANEL_2
        add_rect(s, Inches(0.55), y, Inches(12.2), Inches(0.45),
                 fill=bg, line=None)
        # method pill
        pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(0.85), y + Inches(0.08),
                                   Inches(0.9), Inches(0.3))
        pill.adjustments[0] = 0.5
        pill.fill.solid(); pill.fill.fore_color.rgb = c
        pill.line.fill.background()
        ptf = pill.text_frame
        ptf.margin_left = Emu(0); ptf.margin_right = Emu(0)
        ptf.margin_top = Emu(0); ptf.margin_bottom = Emu(0)
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run()
        pr.text = method; pr.font.name = FONT_BODY
        pr.font.size = Pt(9); pr.font.bold = True; pr.font.color.rgb = BG
        add_text_box(s, Inches(2.0), y + Inches(0.1),
                     Inches(4.4), Inches(0.4),
                     path, size=11, color=TEXT, font=FONT_MONO)
        add_text_box(s, Inches(6.5), y + Inches(0.1),
                     Inches(6.1), Inches(0.4),
                     desc, size=11, color=DIM)
        y += Inches(0.48)
def color_for_method(_):
    return DIM
slide_api(); slides_def.append("api")

# ============ Slide: 測試 ============
def slide_testing():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_slide_header(s, "06 · Verification", "驗證與測試結果")
    rows = [
        ("前端 TS 型別檢查",           "PASS", GREEN, "tsc --noEmit -p tsconfig.app.json"),
        ("後端 TS 型別檢查",           "PASS", GREEN, "tsc --noEmit -p tsconfig.server.json"),
        ("Production build",          "PASS", GREEN, "1946 modules · 1.1MB JS · 38KB CSS"),
        ("/api/analyze 假投資訊息",    "PASS", GREEN, "score=96 · level=high · 5 keywords"),
        ("/api/verify-site .xyz",     "PASS", GREEN, "verdict=danger · score=84"),
        ("SSRF 阻擋 127.0.0.1 / 169.254.169.254 / localhost",
                                     "PASS", GREEN, "ssrfBlocked=true · 三種情境皆擋"),
        ("SQLite 重啟資料保留",        "PASS", GREEN, "16 posts + 1 report 完整"),
        ("Neo4j MERGE 增量同步",      "PASS", GREEN, "移除全砍 + 7 天 stale 清理"),
        ("LLM Markdown 輸出渲染",     "PASS", GREEN, "標題/條列/表格/連結 完整"),
        ("KG 熱載過濾(無刷新)",      "PASS", GREEN, "chip 切換即時 + 700ms 重排"),
        ("165 案例去重(URL hash)",   "PASS", GREEN, "改 case-detail/{Id} 避免折疊"),
        ("政治 / PSA 噪聲過濾",       "PASS", GREEN, "BLOCK/ALLOW 去除 ~20+ 噪聲"),
    ]
    add_rect(s, Inches(0.55), Inches(1.55),
             Inches(12.2), Inches(0.5), fill=PANEL_2)
    add_text_box(s, Inches(0.85), Inches(1.62), Inches(6),
                 Inches(0.4), "測試項目", size=11, color=DIM, bold=True)
    add_text_box(s, Inches(7.4), Inches(1.62), Inches(1.2),
                 Inches(0.4), "結果", size=11, color=DIM, bold=True)
    add_text_box(s, Inches(8.7), Inches(1.62), Inches(4),
                 Inches(0.4), "備註", size=11, color=DIM, bold=True)
    y = Inches(2.1)
    for i, (t, st, c, note) in enumerate(rows):
        bg = PANEL if i % 2 == 0 else PANEL_2
        add_rect(s, Inches(0.55), y, Inches(12.2),
                 Inches(0.4), fill=bg, line=None)
        add_text_box(s, Inches(0.85), y + Inches(0.07),
                     Inches(6.5), Inches(0.35),
                     t, size=10, color=TEXT)
        pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(7.4), y + Inches(0.07),
                                   Inches(0.85), Inches(0.25))
        pill.adjustments[0] = 0.5
        pill.fill.solid(); pill.fill.fore_color.rgb = c
        pill.line.fill.background()
        ptf = pill.text_frame
        ptf.margin_left = Emu(0); ptf.margin_right = Emu(0)
        ptf.margin_top = Emu(0); ptf.margin_bottom = Emu(0)
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run()
        pr.text = st; pr.font.name = FONT_BODY
        pr.font.size = Pt(8); pr.font.bold = True; pr.font.color.rgb = BG
        add_text_box(s, Inches(8.7), y + Inches(0.07),
                     Inches(4), Inches(0.35),
                     note, size=10, color=DIM)
        y += Inches(0.42)
slide_testing(); slides_def.append("testing")

# ============ Slide: Demo / Future ============
def slide_demo():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_slide_header(s, "07 · Demo & Roadmap", "現場展示 + 未來工作")
    add_rect(s, Inches(0.55), Inches(1.6),
             Inches(6.0), Inches(5.4), fill=PANEL)
    add_text_box(s, Inches(0.8), Inches(1.75),
                 Inches(5.5), Inches(0.4),
                 "Demo 流程", size=14, color=ACCENT_2, bold=True)
    demo = (
        "1. npm run dev:full (Docker + 前後端一鍵啟動)\n"
        "2. /  Dashboard:看 4 個來源同步狀態\n"
        "3. /cases:分頁瀏覽,確認多來源混合\n"
        "4. /graph:全螢幕 → 點 post 節點 → 按「前往原始來源」\n"
        "5. /verify:輸入 127.0.0.1 驗 SSRF\n"
        "6. /report:送 LINE 投資詐騙訊息 → 自動分析\n"
        "7. 右下浮動 ChatWidget:自然語言對話 → LLM Markdown 回覆"
    )
    add_text_box(s, Inches(0.8), Inches(2.25),
                 Inches(5.5), Inches(4.7),
                 demo, size=11, color=TEXT, line_spacing=1.7)

    add_rect(s, Inches(6.75), Inches(1.6),
             Inches(6.0), Inches(5.4), fill=PANEL)
    add_text_box(s, Inches(7.0), Inches(1.75),
                 Inches(5.5), Inches(0.4),
                 "未來工作(Roadmap)", size=14, color=ACCENT_2, bold=True)
    items = [
        ("P1", "vitest 單元測試覆蓋 analyzeMessage / safeFetch", GREEN),
        ("P1", "規則 / 關鍵字外部化為 JSON config", GREEN),
        ("P1", "結構化 logging (pino) + latency 觀測", GREEN),
        ("P2", "node-cron 定時增量爬,而非僅啟動時跑", YELLOW),
        ("P2", "Bundle 拆分 (vis-network / markdown lazy)", YELLOW),
        ("P2", "Anthropic prompt caching 降 LLM 成本", YELLOW),
        ("P3", "使用者回報 moderation queue", ORANGE),
        ("P3", "Grafana + Prometheus 觀測儀表板", ORANGE),
    ]
    yy = Inches(2.3)
    for tag, t, c in items:
        pill = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(7.0), yy + Inches(0.05),
                                   Inches(0.55), Inches(0.3))
        pill.adjustments[0] = 0.5
        pill.fill.solid(); pill.fill.fore_color.rgb = c
        pill.line.fill.background()
        ptf = pill.text_frame
        ptf.margin_left = Emu(0); ptf.margin_right = Emu(0)
        ptf.margin_top = Emu(0); ptf.margin_bottom = Emu(0)
        ptf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pp = ptf.paragraphs[0]; pp.alignment = PP_ALIGN.CENTER
        pr = pp.add_run()
        pr.text = tag; pr.font.name = FONT_BODY
        pr.font.size = Pt(10); pr.font.bold = True; pr.font.color.rgb = BG
        add_text_box(s, Inches(7.7), yy + Inches(0.07),
                     Inches(5.0), Inches(0.35),
                     t, size=11, color=TEXT)
        yy += Inches(0.55)
slide_demo(); slides_def.append("demo")

# ============ Slide: 參考資料 ============
def slide_references():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_slide_header(s, "08 · References", "參考資料")
    cats = [
        ("資料來源", ACCENT_2, [
            "165 全民防騙網 — https://165dashboard.tw",
            "PTT Web 介面 — https://www.ptt.cc",
            "Google News RSS — news.google.com/rss/search",
            "OpenPhish 公開 feed — https://openphish.com/feed.txt",
            "PhishTank — https://phishtank.org",
            "Google Safe Browsing v4 — developers.google.com/safe-browsing",
        ]),
        ("套件 / 框架", GREEN, [
            "React 19 — react.dev",
            "Vite 7 — vitejs.dev",
            "Express 5 — expressjs.com",
            "Neo4j Driver — neo4j.com/docs/javascript-manual",
            "vis-network — visjs.github.io/vis-network",
            "better-sqlite3 — github.com/WiseLibs/better-sqlite3",
            "cheerio — cheerio.js.org",
            "react-markdown / remark-gfm",
        ]),
        ("模型 / API", ORANGE, [
            "Ollama — ollama.com",
            "qwen3.5:9b (對話 + JSON 模式)",
            "multilingual-e5-large — huggingface.co/intfloat/multilingual-e5-large",
            "OpenAI Chat Completions 規格",
        ]),
        ("學術 / 文獻", PURPLE, [
            "Reimers & Gurevych (2019) Sentence-BERT",
            "Wang et al. (2024) Multilingual E5 Text Embeddings",
            "OWASP Top 10 — SSRF (A10:2021)",
            "OWASP API Security Top 10",
            "MITRE ATT&CK — Phishing (T1566)",
            "刑事警察局 165 反詐騙年度報告",
        ]),
    ]
    cw = Inches(6.0)
    gap = Inches(0.25)
    sx = Inches(0.55)
    y_start = Inches(1.6)
    for i, (title, c, items) in enumerate(cats):
        col = i % 2
        row = i // 2
        x = sx + col * (cw + gap)
        y = y_start + row * Inches(2.85)
        add_rect(s, x, y, cw, Inches(2.7), fill=PANEL)
        add_text_box(s, x + Inches(0.3), y + Inches(0.2),
                     cw - Inches(0.6), Inches(0.4),
                     title, size=14, color=c, bold=True)
        body = "\n".join("• " + it for it in items)
        add_text_box(s, x + Inches(0.3), y + Inches(0.7),
                     cw - Inches(0.6), Inches(2.0),
                     body, size=9.5, color=TEXT, line_spacing=1.5)
slide_references(); slides_def.append("references")

# ============ Slide: Q & A ============
def slide_qa():
    s = prs.slides.add_slide(BLANK)
    set_bg(s)
    add_text_box(s, Inches(0.55), Inches(2.5),
                 Inches(12.2), Inches(1.5),
                 "Q & A",
                 size=72, color=TEXT, bold=True, font=FONT_TITLE,
                 align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.55), Inches(4.0),
                 Inches(12.2), Inches(0.5),
                 "謝謝聆聽",
                 size=20, color=ACCENT_2,
                 align=PP_ALIGN.CENTER)
    add_text_box(s, Inches(0.55), Inches(4.7),
                 Inches(12.2), Inches(0.5),
                 "github.com/jason79461385/Crawler_KG_ANTI_FAKE",
                 size=12, color=MUTED, align=PP_ALIGN.CENTER)
slide_qa(); slides_def.append("qa")

# ============ 為所有 slide 加 footer ============
total = len(prs.slides)
for i, slide in enumerate(prs.slides, start=1):
    # skip cover / qa
    if i == 1 or i == total:
        continue
    add_footer(slide, i, total)

out = "/Users/jkbry/Documents/New project/Scam_Intel_Console_Slides.pptx"
prs.save(out)
print(f"✓ 已輸出 {total} 頁:{out}")
